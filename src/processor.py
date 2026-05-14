"""
Document processor: handles PDFs (native + OCR fallback), DOCX, PPTX/PPT, images, and plain text.
Extracts clean text and structured legal fields using Claude.
"""
import json
import logging
import re
import uuid
from pathlib import Path
from typing import Optional

import pdfplumber
from anthropic import Anthropic

from .config import ANTHROPIC_API_KEY, CLAUDE_MODEL, CHUNK_SIZE, CHUNK_OVERLAP
from .models import ProcessedDocument

logger = logging.getLogger(__name__)
_client = Anthropic(api_key=ANTHROPIC_API_KEY)

_STRUCTURED_PROMPT = """You are a legal document analyst. Extract structured information from the text below.

Return ONLY a valid JSON object with these fields (use null for missing fields):
{{
  "document_type": "complaint/contract/notice/memo/order/other",
  "case_title": "case or matter title if present",
  "case_number": "docket or case number if present",
  "parties": [{{"name": "...", "role": "plaintiff/defendant/counsel/other"}}],
  "dates": [{{"date": "YYYY-MM-DD or as written", "event": "what happened"}}],
  "legal_issues": ["identified legal claims or issues"],
  "key_facts": ["important factual statements"],
  "referenced_statutes": ["laws, codes, or regulations cited"],
  "jurisdiction": "court or jurisdiction if mentioned",
  "filing_status": "filed/pending/dismissed/other if mentioned",
  "document_quality": "note any OCR artifacts, illegible sections, or gaps"
}}

Document text:
{text}"""


def _clean_text(text: str) -> str:
    """Normalize whitespace and remove common OCR artifacts."""
    # Common OCR substitutions
    text = re.sub(r'\|(?=[A-Z])', 'I', text)
    # Remove non-printable chars except standard whitespace
    text = re.sub(r'[^\x09\x0A\x0D\x20-\x7E]', ' ', text)
    # Collapse excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    return text.strip()


def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list:
    """Split text into overlapping paragraph-aware chunks."""
    paragraphs = [p.strip() for p in re.split(r'\n\n+', text) if p.strip()]
    chunks = []
    current = ""

    for para in paragraphs:
        if len(current) + len(para) + 2 <= chunk_size:
            current = (current + "\n\n" + para).strip()
        else:
            if current:
                chunks.append(current)
                # carry overlap from tail of current chunk
                words = current.split()
                tail = " ".join(words[-(overlap // 5):]) if len(words) > overlap // 5 else ""
                current = (tail + "\n\n" + para).strip()
            else:
                # paragraph itself exceeds chunk_size — split by sentence
                for sent in re.split(r'(?<=[.!?])\s+', para):
                    if len(current) + len(sent) + 1 <= chunk_size:
                        current = (current + " " + sent).strip()
                    else:
                        if current:
                            chunks.append(current)
                        current = sent

    if current:
        chunks.append(current)

    return [c for c in chunks if len(c.strip()) > 40]


def _pdf_native(file_path: str) -> tuple:
    """Extract text from a PDF using pdfplumber."""
    pages, notes = [], []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            txt = page.extract_text() or ""
            if len(txt.strip()) < 20:
                notes.append(f"Page {i+1}: low native text — may be scanned")
            pages.append(txt)
    return "\n\n".join(pages), notes


def _pdf_ocr(file_path: str) -> tuple:
    """OCR a scanned PDF using pdf2image + pytesseract."""
    notes = []
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(file_path, dpi=300)
        texts = []
        for i, img in enumerate(images):
            txt = pytesseract.image_to_string(img, config="--oem 3 --psm 6")
            texts.append(txt)
            if len(txt.strip()) < 50:
                notes.append(f"Page {i+1}: OCR returned sparse text")
        return "\n\n".join(texts), notes

    except ImportError:
        notes.append("OCR dependencies unavailable (pdf2image / pytesseract). Install them for scanned-PDF support.")
        return "", notes
    except Exception as exc:
        notes.append(f"OCR failed: {exc}")
        return "", notes


def _image_ocr(file_path: str) -> str:
    """OCR a single image file."""
    import pytesseract
    from PIL import Image
    return pytesseract.image_to_string(Image.open(file_path), config="--oem 3 --psm 6")


def _docx_extract(file_path: str) -> tuple:
    """Extract text from a .docx file using python-docx."""
    notes = []
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        text = "\n\n".join(parts)
        if not text.strip():
            notes.append("DOCX contained no extractable text paragraphs")
        return text, notes
    except ImportError:
        raise RuntimeError("python-docx required for .docx files. Run: pip install python-docx")
    except Exception as exc:
        raise ValueError(f"Failed to read .docx: {exc}")


def _pptx_extract(file_path: str) -> tuple:
    """Extract text from a .pptx file using python-pptx."""
    notes = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
            else:
                notes.append(f"Slide {slide_num}: no text found")
        text = "\n\n".join(parts)
        if not text.strip():
            notes.append("PPTX contained no extractable text")
        return text, notes
    except ImportError:
        raise RuntimeError("python-pptx required for .pptx/.ppt files. Run: pip install python-pptx")
    except Exception as exc:
        raise ValueError(f"Failed to read presentation: {exc}")


def _extract_structured(text: str) -> dict:
    """Use Claude to pull structured legal fields from raw text."""
    truncated = text[:8000]
    response = _client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=2000,
        messages=[{"role": "user", "content": _STRUCTURED_PROMPT.format(text=truncated)}],
    )
    raw = response.content[0].text.strip()
    # Strip markdown fences if present
    raw = re.sub(r'^```(?:json)?\n?', '', raw)
    raw = re.sub(r'\n?```$', '', raw)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return {"raw_extraction": raw, "parse_error": "Response was not valid JSON"}


def process_document(file_path: str, doc_id: Optional[str] = None) -> ProcessedDocument:
    """
    Main entry point.  Accepts .pdf, .txt, .png/.jpg/.tiff images.
    Returns a ProcessedDocument ready for retrieval and drafting.
    """
    if doc_id is None:
        doc_id = str(uuid.uuid4())[:8]

    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {file_path}")

    suffix = path.suffix.lower()
    notes: list = []
    extraction_method = "unknown"

    # ── Text extraction ───────────────────────────────────────────
    if suffix == ".txt":
        raw_text = path.read_text(encoding="utf-8", errors="replace")
        extraction_method = "text"

    elif suffix == ".pdf":
        native_text, native_notes = _pdf_native(file_path)
        notes.extend(native_notes)

        scanned_page_ratio = len([n for n in native_notes if "low native" in n]) / max(1, len(native_notes) + 1)
        if len(native_text.strip()) < 200 or scanned_page_ratio > 0.5:
            notes.append("Native extraction insufficient — attempting OCR")
            ocr_text, ocr_notes = _pdf_ocr(file_path)
            notes.extend(ocr_notes)
            if len(ocr_text) > len(native_text):
                raw_text = ocr_text
                extraction_method = "ocr"
            else:
                raw_text = native_text
                extraction_method = "native_pdf_limited"
                notes.append("OCR did not improve on native extraction; using native text")
        else:
            raw_text = native_text
            extraction_method = "native_pdf"

    elif suffix == ".docx":
        raw_text, docx_notes = _docx_extract(file_path)
        notes.extend(docx_notes)
        extraction_method = "docx"

    elif suffix in (".pptx", ".ppt"):
        raw_text, pptx_notes = _pptx_extract(file_path)
        notes.extend(pptx_notes)
        extraction_method = "pptx"

    elif suffix in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
        try:
            raw_text = _image_ocr(file_path)
            extraction_method = "ocr_image"
        except ImportError:
            raise RuntimeError("pytesseract + Pillow required for image files. Install them first.")

    else:
        try:
            raw_text = path.read_text(encoding="utf-8", errors="replace")
            extraction_method = "text_fallback"
            notes.append(f"Unknown extension '{suffix}' — treated as plain text")
        except Exception as exc:
            raise ValueError(f"Cannot process file type '{suffix}': {exc}")

    raw_text = _clean_text(raw_text)
    if not raw_text:
        notes.append("WARNING: extracted text is empty after cleaning")
        raw_text = "[No extractable text found in this document]"

    # ── Structured field extraction ───────────────────────────────
    logger.info(f"[{doc_id}] Extracting structured fields via Claude...")
    structured_fields = _extract_structured(raw_text)

    # ── Chunking ──────────────────────────────────────────────────
    chunks = _chunk_text(raw_text)
    notes.append(f"Produced {len(chunks)} chunks (method={extraction_method})")

    return ProcessedDocument(
        doc_id=doc_id,
        file_path=str(file_path),
        raw_text=raw_text,
        structured_fields=structured_fields,
        chunks=chunks,
        extraction_method=extraction_method,
        processing_notes=notes,
    )
